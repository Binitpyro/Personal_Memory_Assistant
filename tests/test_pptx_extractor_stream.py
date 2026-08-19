"""The OOXML-reading PPTX extractor must produce exactly what python-pptx did.

`PptxExtractor` stopped using python-pptx because `Presentation()` builds an
lxml element tree for every part in the package before a single character is
read. Measured over 42 real decks (33 MB, of which 0.1 MB is media and 0.15 MB
is slide XML): **90.0 MB** peak working set through python-pptx against
**2.2 MB** reading the same text out of the zip. Through the real pipeline the
PPTX-only ingestion floor fell from 111.9 MB to 30.2 MB. See CLAUDE.md section 6.

Hand-rolled XML replacing a library is exactly the change that loses content
quietly, and the tests that existed before mocked python-pptx out entirely, so
they could not have caught it. These compare the two implementations on decks
built to exercise every branch, and the reference implementation below is a
verbatim copy of the code that was replaced.

Run against 42 real lecture decks during development: 42/42 identical, after the
comparison caught a genuine defect - the first version read every text frame on
a notes part, so the slide-number placeholder was appended to every note.
"""

from pathlib import Path

import pytest

from app.indexing.extractors.pptx_extractor import PptxExtractor

MAX_SIZE = 50_000_000

pytest.importorskip("pptx", reason="python-pptx is the reference implementation here")


def _reference_extract(path: Path, max_file_size: int) -> list[str]:
    """The python-pptx implementation this replaced, kept verbatim.

    Deliberately not refactored or tidied: its value is being the thing that
    actually shipped, including the `hasattr(shape, "text")` test that silently
    skips pictures, connectors and group shapes.
    """
    from pptx import Presentation

    out: list[str] = []
    prs = Presentation(str(path))
    total = 0
    for i, slide in enumerate(prs.slides):
        out.append(f"--- Slide {i + 1} ---")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
                total += len(shape.text)

        for shape in slide.shapes:
            try:
                if not shape.has_table:
                    continue
            except Exception:
                continue
            try:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        out.append(row_text)
                        total += len(row_text)
            except Exception:
                pass

        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text and notes_text.strip():
                    out.append(f"[Notes] {notes_text}")
                    total += len(notes_text)
        except Exception:
            pass

        if total > max_file_size:
            break
    return out


def _build_deck(path: Path) -> Path:
    """A deck exercising titles, body text, multi-paragraph frames, tables,
    notes, a picture-free group shape, and slides with nothing on them."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank, title_only = prs.slide_layouts[6], prs.slide_layouts[5]

    s1 = prs.slides.add_slide(title_only)
    s1.shapes.title.text = "Quarterly Review"
    box = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(2))
    frame = box.text_frame
    frame.text = "First paragraph"
    frame.add_paragraph().text = "Second paragraph"
    frame.add_paragraph().text = ""
    s1.notes_slide.notes_text_frame.text = "Speaker notes with detail"

    s2 = prs.slides.add_slide(title_only)
    s2.shapes.title.text = "Numbers"
    table = s2.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    for r, row in enumerate((("Region", "Revenue"), ("EMEA", "12"), ("", ""))):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # A slide with no shapes at all: the marker must still be emitted.
    prs.slides.add_slide(blank)

    s4 = prs.slides.add_slide(blank)
    tb = s4.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "Run one "
    tb.text_frame.paragraphs[0].add_run().text = "and run two"
    tb.text_frame.paragraphs[0].font.size = Pt(18)

    prs.save(str(path))
    return path


class TestEquivalenceWithPythonPptx:
    def test_output_is_identical_on_a_multi_feature_deck(self, tmp_path):
        deck = _build_deck(tmp_path / "features.pptx")

        assert list(PptxExtractor().extract_stream(deck, MAX_SIZE)) == _reference_extract(
            deck, MAX_SIZE
        )

    def test_output_is_identical_when_there_are_no_notes(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "No notes here"
        deck = tmp_path / "nonotes.pptx"
        prs.save(str(deck))

        assert list(PptxExtractor().extract_stream(deck, MAX_SIZE)) == _reference_extract(
            deck, MAX_SIZE
        )

    def test_output_is_identical_on_an_empty_deck(self, tmp_path):
        from pptx import Presentation

        deck = tmp_path / "empty.pptx"
        Presentation().save(str(deck))

        assert list(PptxExtractor().extract_stream(deck, MAX_SIZE)) == _reference_extract(
            deck, MAX_SIZE
        )


class TestStreamingProperties:
    def test_slide_order_follows_the_presentation_not_the_filenames(self, tmp_path):
        """Slide part numbering does not have to match display order - a deck
        that has had slides reordered or deleted will not sort correctly by
        filename, so order comes from p:sldIdLst through the rels."""
        from pptx import Presentation

        prs = Presentation()
        for name in ("Alpha", "Bravo", "Charlie"):
            prs.slides.add_slide(prs.slide_layouts[5]).shapes.title.text = name
        deck = tmp_path / "order.pptx"
        prs.save(str(deck))

        # Reverse the slide id list in place, so display order no longer
        # matches slide1/slide2/slide3.
        import shutil
        import zipfile

        rewritten = tmp_path / "reordered.pptx"
        with zipfile.ZipFile(deck) as src:
            names = src.namelist()
            data = {n: src.read(n) for n in names}
        pres = data["ppt/presentation.xml"].decode("utf-8")
        start, end = pres.index("<p:sldIdLst>"), pres.index("</p:sldIdLst>")
        inner = pres[start + len("<p:sldIdLst>") : end]
        entries = [e + "/>" for e in inner.split("/>") if e.strip()]
        data["ppt/presentation.xml"] = (
            pres[: start + len("<p:sldIdLst>")] + "".join(reversed(entries)) + pres[end:]
        ).encode("utf-8")
        with zipfile.ZipFile(rewritten, "w", zipfile.ZIP_DEFLATED) as dst:
            for n in names:
                dst.writestr(n, data[n])
        shutil.copystat(deck, rewritten)

        titles = [
            ln
            for ln in PptxExtractor().extract_stream(rewritten, MAX_SIZE)
            if ln in ("Alpha", "Bravo", "Charlie")
        ]
        assert titles == ["Charlie", "Bravo", "Alpha"], (
            "slide order was taken from part filenames rather than p:sldIdLst"
        )

    def test_python_pptx_is_never_imported(self, tmp_path):
        """The whole point of the rewrite. If this fails the memory win is gone,
        because Presentation() builds the lxml tree the extractor exists to
        avoid."""
        import sys

        deck = _build_deck(tmp_path / "noimport.pptx")
        for mod in [m for m in sys.modules if m == "pptx" or m.startswith("pptx.")]:
            del sys.modules[mod]

        list(PptxExtractor().extract_stream(deck, MAX_SIZE))

        assert not [m for m in sys.modules if m == "pptx" or m.startswith("pptx.")], (
            "the extractor imported python-pptx"
        )

    def test_a_corrupt_slide_part_does_not_lose_the_rest_of_the_deck(self, tmp_path):
        import zipfile

        deck = _build_deck(tmp_path / "corrupt-part.pptx")
        broken = tmp_path / "broken.pptx"
        with zipfile.ZipFile(deck) as src:
            names = src.namelist()
            data = {n: src.read(n) for n in names}
        target = next(n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
        data[target] = b"<p:sld><not-closed>"
        with zipfile.ZipFile(broken, "w", zipfile.ZIP_DEFLATED) as dst:
            for n in names:
                dst.writestr(n, data[n])

        out = list(PptxExtractor().extract_stream(broken, MAX_SIZE))
        assert any("Numbers" in ln or "Quarterly Review" in ln for ln in out), (
            "one unparseable slide discarded the whole deck"
        )


class TestXmlEntityExpansion:
    """A deck is untrusted input: PMA indexes documents the user did not write.

    Both ElementTree and lxml expand internal entities, so a "billion laughs"
    prolog inflates during parse - measured at 30,000 characters from a 9-line
    declaration at depth 4, scaling geometrically with depth. python-pptx used
    lxml, so this exposure predates the rewrite rather than arriving with it.
    `_ooxml_guard` bounds the *declared* decompressed size and cannot see
    post-parse expansion.

    OOXML never legitimately declares a DTD, so refusing one closes the vector
    without a defusedxml dependency.
    """

    def _repack(self, tmp_path, deck, part, payload):
        import zipfile

        out = tmp_path / "bomb.pptx"
        with zipfile.ZipFile(deck) as src:
            names = src.namelist()
            data = {n: src.read(n) for n in names}
        data[part] = payload
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
            for n in names:
                dst.writestr(n, data[n])
        return out

    def test_a_slide_declaring_a_dtd_is_refused(self, tmp_path):
        deck = _build_deck(tmp_path / "src.pptx")
        bomb = (
            b'<?xml version="1.0"?>\n<!DOCTYPE lolz [\n'
            b' <!ENTITY lol "lol">\n'
            b' <!ENTITY lol1 "&lol;&lol;&lol;&lol;&lol;&lol;&lol;&lol;&lol;&lol;">\n'
            b' <!ENTITY lol2 "&lol1;&lol1;&lol1;&lol1;&lol1;&lol1;&lol1;&lol1;&lol1;&lol1;">\n'
            b"]>\n<r>&lol2;</r>"
        )
        import zipfile

        with zipfile.ZipFile(deck) as z:
            target = next(
                n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
        packed = self._repack(tmp_path, deck, target, bomb)

        out = list(PptxExtractor().extract_stream(packed, MAX_SIZE))

        assert not any("lollol" in line for line in out), "entity expansion reached the output"
        # The rest of the deck must survive one refused part.
        assert any(line.startswith("--- Slide") for line in out)

    def test_the_guard_is_what_refuses_it(self, tmp_path):
        """Direct, so a future refactor that keeps output clean by accident but
        drops the check still fails."""
        import zipfile

        from app.indexing.extractors.pptx_extractor import _parse_part

        deck = _build_deck(tmp_path / "direct.pptx")
        packed = self._repack(
            tmp_path, deck, "ppt/presentation.xml", b'<!DOCTYPE x [<!ENTITY a "b">]><x/>'
        )
        with (
            zipfile.ZipFile(packed) as z,
            pytest.raises(ValueError, match="DTD"),
        ):
            _parse_part(z, "ppt/presentation.xml")
