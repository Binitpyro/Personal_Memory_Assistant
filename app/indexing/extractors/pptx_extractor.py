import logging
from collections.abc import Iterator
from pathlib import Path

logger = logging.getLogger(__name__)


class PptxExtractor:
    def can_handle(self, path: Path) -> bool:
        return path.suffix.lower() == ".pptx"

    def extract_stream(self, path: Path, max_file_size: int) -> Iterator[str]:
        """Yield text from slides in a PPTX document.

        Order per slide is: marker, regular shape text, tables, speaker
        notes. Regular text must stay first - app/indexing/summarizer.py's
        _summarize_doc_text() takes the first line after each "--- Slide
        N ---" marker as that slide's title, and tables/notes are appended
        content, not the title.
        """
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(str(path))
            total = 0
            for i, slide in enumerate(prs.slides):
                yield f"--- Slide {i + 1} ---"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        yield shape.text
                        total += len(shape.text)

                for shape in slide.shapes:
                    try:
                        if not shape.has_table:
                            continue
                    except Exception:
                        continue
                    try:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells)
                            if row_text.strip():
                                yield row_text
                                total += len(row_text)
                    except Exception as table_err:
                        logger.debug(
                            "Skipping unreadable table on slide %d of %s: %s",
                            i + 1,
                            path,
                            table_err,
                        )

                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text and notes_text.strip():
                            yield f"[Notes] {notes_text}"
                            total += len(notes_text)
                except Exception as notes_err:
                    logger.debug(
                        "Skipping unreadable notes on slide %d of %s: %s", i + 1, path, notes_err
                    )

                if total > max_file_size:
                    break
        except Exception as e:
            logger.warning("Failed to extract PPTX %s: %s", path, e)

    def extract(self, path: Path, max_file_size: int) -> str:
        """Legacy extraction for backward compatibility."""
        return "\n".join(self.extract_stream(path, max_file_size))[:max_file_size]
