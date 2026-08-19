"""
tests/test_extractors.py
Covers all 7 file extractors in app/indexing/extractors/:
- CsvExtractor, JsonExtractor, PdfExtractor (mocked)
- DocxExtractor, XlsxExtractor, PptxExtractor, EpubExtractor (mocked)
"""

import json
import zipfile
from unittest.mock import MagicMock, patch

from app.config import settings
from app.indexing.extractors import EXTRACTORS, ExtractMeta
from app.indexing.extractors.csv_extractor import CsvExtractor
from app.indexing.extractors.docx_extractor import DocxExtractor
from app.indexing.extractors.epub_extractor import EpubExtractor
from app.indexing.extractors.json_extractor import JsonExtractor
from app.indexing.extractors.pdf_extractor import PdfExtractor
from app.indexing.extractors.pptx_extractor import PptxExtractor
from app.indexing.extractors.xlsx_extractor import XlsxExtractor

MAX_SIZE = 1_000_000


# ── CSV ───────────────────────────────────────────────────────────────────────


class TestCsvExtractor:
    def setup_method(self):
        self.ext = CsvExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "data.csv")
        assert not self.ext.can_handle(tmp_path / "data.txt")

    def test_basic_csv(self, tmp_path):
        f = tmp_path / "test.csv"
        f.write_text("name,age\nAlice,30\nBob,25\n", encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert "Alice" in result
        assert "name: Alice" in result or "name" in result

    def test_empty_csv(self, tmp_path):
        f = tmp_path / "empty.csv"
        f.write_text("", encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert result == ""

    def test_headers_only_csv(self, tmp_path):
        f = tmp_path / "headers_only.csv"
        f.write_text("col1,col2,col3\n", encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert result == ""

    def test_max_size_truncation(self, tmp_path):
        f = tmp_path / "big.csv"
        rows = ["a,b"] + [f"value{i},{i}" for i in range(100)]
        f.write_text("\n".join(rows), encoding="utf-8")
        result = self.ext.extract(f, 50)
        assert len(result) <= 50

    def test_corrupt_file_returns_empty(self, tmp_path):
        # File that doesn't exist
        result = self.ext.extract(tmp_path / "nonexistent.csv", MAX_SIZE)
        assert result == ""

    def test_rows_limit(self, tmp_path):
        """Should complete without error on a file past the row floor."""
        f = tmp_path / "large.csv"
        lines = ["col"] + [str(i) for i in range(6000)]
        f.write_text("\n".join(lines), encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert isinstance(result, str)

    def test_row_cap_scales_with_the_size_budget(self):
        """A flat cap truncated narrow-row files long before they reached the
        byte budget the caller actually set. The cap now tracks that budget,
        with a floor for small budgets and a ceiling for pathological input."""
        from app.indexing.extractors.csv_extractor import (
            _MAX_CSV_ROWS,
            _MIN_CSV_ROWS,
            _row_budget,
        )

        assert _row_budget(1_000) == _MIN_CSV_ROWS  # floor
        assert _row_budget(10**12) == _MAX_CSV_ROWS  # ceiling
        assert _MIN_CSV_ROWS < _row_budget(50 * 1024 * 1024) <= _MAX_CSV_ROWS
        # Monotonic: a bigger budget never yields a smaller cap.
        assert _row_budget(10 * 1024 * 1024) <= _row_budget(50 * 1024 * 1024)

    def test_rows_limit_warns_on_truncation(self, tmp_path, caplog):
        """P0-3: the row cap used to truncate silently. It must log a warning so
        a truncated import is diagnosable."""
        from app.indexing.extractors.csv_extractor import _MIN_CSV_ROWS

        # 200_000 pins the row cap at the floor (200_000 // 40 == 5000) while
        # leaving far more byte budget than ~45KB of narrow rows consumes, so
        # the row cap is the binding constraint rather than the byte budget.
        budget = 200_000
        f = tmp_path / "large.csv"
        lines = ["col"] + [str(i) for i in range(_MIN_CSV_ROWS + 500)]
        f.write_text("\n".join(lines), encoding="utf-8")

        with caplog.at_level("WARNING", logger="app.indexing.extractors.csv_extractor"):
            rows = list(self.ext.extract_stream(f, budget))

        assert len(rows) == _MIN_CSV_ROWS
        assert any("exceeds" in r.message and "truncating" in r.message for r in caplog.records)

    def test_byte_budget_still_wins_when_it_is_tighter(self, tmp_path):
        """Scaling the row cap must not let a wide-row file blow past the byte
        budget - whichever limit binds first still stops extraction."""
        f = tmp_path / "wide.csv"
        header = "a,b,c"
        rows = [",".join(["x" * 200] * 3) for _ in range(500)]
        f.write_text("\n".join([header, *rows]), encoding="utf-8")

        out = list(self.ext.extract_stream(f, 5_000))

        assert len(out) < 500, "byte budget did not stop a wide-row file"
        assert sum(len(r) for r in out) <= 5_000 + len(out[-1])

    def test_under_cap_does_not_warn(self, tmp_path, caplog):
        f = tmp_path / "small.csv"
        lines = ["col"] + [str(i) for i in range(10)]
        f.write_text("\n".join(lines), encoding="utf-8")

        with caplog.at_level("WARNING", logger="app.indexing.extractors.csv_extractor"):
            list(self.ext.extract_stream(f, MAX_SIZE))

        assert not any("truncating" in r.message for r in caplog.records)


# ── JSON ──────────────────────────────────────────────────────────────────────


class TestJsonExtractor:
    def setup_method(self):
        self.ext = JsonExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "data.json")
        assert not self.ext.can_handle(tmp_path / "data.csv")

    def test_basic_json_dict(self, tmp_path):
        f = tmp_path / "test.json"
        f.write_text(json.dumps({"name": "Alice", "age": 30}), encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert "Alice" in result or "name" in result

    def test_json_list(self, tmp_path):
        f = tmp_path / "list.json"
        f.write_text(json.dumps([1, 2, 3]), encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert isinstance(result, str)

    def test_invalid_json(self, tmp_path):
        f = tmp_path / "bad.json"
        f.write_text("{not valid json}", encoding="utf-8")
        result = self.ext.extract(f, MAX_SIZE)
        assert isinstance(result, str)

    def test_missing_file(self, tmp_path):
        result = self.ext.extract(tmp_path / "missing.json", MAX_SIZE)
        assert result == ""


# ── PDF (mocked pypdf) ────────────────────────────────────────────────────────


class TestPdfExtractor:
    def setup_method(self):
        self.ext = PdfExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "doc.pdf")
        assert not self.ext.can_handle(tmp_path / "doc.docx")

    def test_extract_success(self, tmp_path):
        fake_path = tmp_path / "test.pdf"
        fake_path.touch()
        mock_pypdf = MagicMock()
        mock_page = MagicMock()
        mock_page.extract_text.return_value = "Hello world from PDF"
        mock_reader = MagicMock()
        mock_reader.is_encrypted = False
        mock_reader.pages = [mock_page]
        mock_pypdf.PdfReader.return_value = mock_reader
        with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
            result = self.ext.extract(fake_path, MAX_SIZE)
        assert "Hello world" in result
        mock_pypdf.PdfReader.assert_called_once_with(str(fake_path), strict=False)

    def test_encrypted_pdf(self, tmp_path):
        fake_path = tmp_path / "encrypted.pdf"
        fake_path.touch()
        mock_pypdf = MagicMock()
        mock_reader = MagicMock()
        mock_reader.is_encrypted = True
        mock_pypdf.PdfReader.return_value = mock_reader
        with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
            result = self.ext.extract(fake_path, MAX_SIZE)
        assert "ENCRYPTED" in result.upper()

    def test_missing_file_returns_empty(self, tmp_path):
        result = self.ext.extract(tmp_path / "missing.pdf", MAX_SIZE)
        assert result == ""

    def test_per_page_exception_skipped(self, tmp_path):
        fake_path = tmp_path / "corrupt_page.pdf"
        fake_path.touch()
        mock_pypdf = MagicMock()
        mock_page1 = MagicMock()
        mock_page1.extract_text.side_effect = OSError("Page 1 corrupt")
        mock_page2 = MagicMock()
        mock_page2.extract_text.return_value = "Page 2 content"
        mock_reader = MagicMock()
        mock_reader.is_encrypted = False
        mock_reader.pages = [mock_page1, mock_page2]
        mock_pypdf.PdfReader.return_value = mock_reader
        with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
            result = self.ext.extract(fake_path, MAX_SIZE)
        assert "Page 2 content" in result


# ── PDF OCR gate integration ─────────────────────────────────────────────────


class TestPdfExtractorOcrGate:
    """The gate runs inside extract_stream and reports via a trailing ExtractMeta."""

    def setup_method(self):
        self.ext = PdfExtractor()

    @staticmethod
    def _page(text, *, scanned=False):
        """A page mock whose gate verdict is controlled by `scanned`."""
        page = MagicMock()
        page.extract_text.return_value = text

        xobjects = {}
        if scanned:
            image = MagicMock()
            image.get_object.return_value = {"/Subtype": "/Image", "/Width": 10, "/Height": 10}
            xobjects["/Im0"] = image

        resources = MagicMock()
        xobj_container = MagicMock()
        xobj_container.get_object.return_value = xobj_container
        xobj_container.__iter__ = lambda self: iter(xobjects)
        xobj_container.raw_get = lambda key: xobjects[key]
        resources.get_object.return_value = resources
        resources.get = lambda key, default=None: (
            xobj_container if key == "/XObject" and xobjects else default
        )

        contents = MagicMock()
        contents.get_object.return_value = contents
        contents.get = lambda key, default=None: 0 if key == "/Length" else default

        page.get = lambda key, default=None: {
            "/Resources": resources,
            "/Contents": contents,
        }.get(key, default)
        return page

    def _run(self, tmp_path, pages, monkeypatch, *, enabled=True, encrypted=False):
        monkeypatch.setattr(settings, "ocr_enabled", enabled)
        monkeypatch.setattr(settings, "ocr_tier", "cpu" if enabled else "none")

        fake_path = tmp_path / "doc.pdf"
        fake_path.touch()
        mock_pypdf = MagicMock()
        mock_reader = MagicMock()
        mock_reader.is_encrypted = encrypted
        mock_reader.pages = pages
        mock_pypdf.PdfReader.return_value = mock_reader
        with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
            return list(self.ext.extract_stream(fake_path, MAX_SIZE))

    def test_disabled_yields_only_strings(self, tmp_path, monkeypatch):
        """With OCR off the stream must be byte-for-byte what it always was."""
        items = self._run(tmp_path, [self._page("Some text " * 30)], monkeypatch, enabled=False)
        assert all(isinstance(i, str) for i in items)
        assert not any(isinstance(i, ExtractMeta) for i in items)

    def test_fully_scanned_pdf_reports_every_page(self, tmp_path, monkeypatch):
        pages = [self._page("", scanned=True) for _ in range(3)]
        items = self._run(tmp_path, pages, monkeypatch)

        meta = items[-1]
        assert isinstance(meta, ExtractMeta)
        assert meta.ocr_pages == (0, 1, 2)
        assert meta.page_count == 3

    def test_mixed_pdf_reports_only_the_scanned_pages(self, tmp_path, monkeypatch):
        pages = [
            self._page("Readable body text that is clearly long enough to pass. " * 4),
            self._page("", scanned=True),
        ]
        items = self._run(tmp_path, pages, monkeypatch)

        meta = items[-1]
        assert meta.ocr_pages == (1,)
        assert meta.native_pages == 1
        # The native page's text is still yielded normally.
        assert any(isinstance(i, str) and "Readable body" in i for i in items)

    def test_meta_is_yielded_exactly_once_and_last(self, tmp_path, monkeypatch):
        pages = [self._page("", scanned=True), self._page("", scanned=True)]
        items = self._run(tmp_path, pages, monkeypatch)

        metas = [i for i in items if isinstance(i, ExtractMeta)]
        assert len(metas) == 1
        assert isinstance(items[-1], ExtractMeta)

    def test_encrypted_pdf_reports_the_reason_and_queues_nothing(self, tmp_path, monkeypatch):
        items = self._run(tmp_path, [], monkeypatch, encrypted=True)

        meta = items[-1]
        assert isinstance(meta, ExtractMeta)
        assert meta.reason == "encrypted"
        assert meta.ocr_pages == ()
        assert "ENCRYPTED" in items[0].upper()

    def test_extract_still_returns_a_string(self, tmp_path, monkeypatch):
        """Regression guard: extract() joins the stream and must skip the meta."""
        monkeypatch.setattr(settings, "ocr_enabled", True)
        monkeypatch.setattr(settings, "ocr_tier", "cpu")

        fake_path = tmp_path / "doc.pdf"
        fake_path.touch()
        mock_pypdf = MagicMock()
        mock_reader = MagicMock()
        mock_reader.is_encrypted = False
        mock_reader.pages = [self._page("Hello from a scanned-ish page " * 5, scanned=True)]
        mock_pypdf.PdfReader.return_value = mock_reader

        with patch.dict("sys.modules", {"pypdf": mock_pypdf}):
            result = self.ext.extract(fake_path, MAX_SIZE)

        assert isinstance(result, str)
        assert "Hello from a scanned-ish page" in result


# ── DOCX (mocked python-docx) ─────────────────────────────────────────────────


class TestDocxExtractor:
    def setup_method(self):
        self.ext = DocxExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "file.docx")
        assert not self.ext.can_handle(tmp_path / "file.pdf")

    def test_extract_success(self, tmp_path):
        fake_path = tmp_path / "test.docx"
        fake_path.touch()
        mock_docx = MagicMock()
        mock_para = MagicMock()
        mock_para.text = "Paragraph text content"

        mock_paragraph_module = MagicMock()
        mock_paragraph_module.Paragraph.return_value = mock_para
        mock_table_module = MagicMock()

        mock_doc = MagicMock()
        mock_child_p = MagicMock()
        mock_child_p.tag = "}p"
        mock_doc.element.body.iterchildren.return_value = [mock_child_p]
        mock_docx.Document.return_value = mock_doc

        with patch.dict(
            "sys.modules",
            {
                "docx": mock_docx,
                "docx.text.paragraph": mock_paragraph_module,
                "docx.table": mock_table_module,
            },
        ):
            result = self.ext.extract(fake_path, MAX_SIZE)
        assert "Paragraph" in result

    def test_missing_file_returns_empty(self, tmp_path):
        result = self.ext.extract(tmp_path / "missing.docx", MAX_SIZE)
        assert result == ""

    def test_docx_heading_hierarchy_is_preserved(self, tmp_path):
        """DOCX carries structure in paragraph styles. Dropping it made a
        section title indistinguishable from a sentence before chunking; the
        extractor now emits markdown markers, which is the same structural
        signal the chunker and the markdown summarizer already consume."""
        from docx import Document

        doc = Document()
        doc.add_heading("Top Level Title", level=0)  # style "Title"
        doc.add_heading("Chapter One", level=1)
        doc.add_paragraph("Ordinary body text under chapter one.")
        doc.add_heading("Section One Point One", level=2)
        doc.add_paragraph("More body text.")
        doc.add_heading("Deep Subsection", level=3)

        path = tmp_path / "structured.docx"
        doc.save(str(path))

        lines = list(self.ext.extract_stream(path, MAX_SIZE))

        assert "# Top Level Title" in lines
        assert "# Chapter One" in lines
        assert "## Section One Point One" in lines
        assert "### Deep Subsection" in lines
        # Body text must stay unprefixed, or everything reads as a heading.
        assert "Ordinary body text under chapter one." in lines
        assert "More body text." in lines

    def test_docx_heading_prefix_survives_a_broken_style(self):
        """Malformed documents can reference a style missing from styles.xml.
        That must degrade to body text, not abort the whole extraction."""
        from unittest.mock import PropertyMock

        from app.indexing.extractors.docx_extractor import _heading_prefix

        broken = MagicMock()
        type(broken).style = PropertyMock(side_effect=KeyError("no such style"))
        assert _heading_prefix(broken) == ""

        styleless = MagicMock()
        styleless.style = None
        assert _heading_prefix(styleless) == ""

        unnamed = MagicMock()
        unnamed.style.name = None
        assert _heading_prefix(unnamed) == ""

        # Word allows Heading 1-9; markdown stops at 6.
        deep = MagicMock()
        deep.style.name = "Heading 9"
        assert _heading_prefix(deep) == "###### "

    def test_docx_extracts_headers_footers_and_footnotes(self, tmp_path):
        """Body text comes first; headers/footers/footnotes are appended,
        each labeled. A linked (inherited) second-section header must not
        duplicate the first section's header text."""
        from docx import Document
        from docx.enum.section import WD_SECTION

        doc = Document()
        doc.add_paragraph("Body paragraph text")

        section1 = doc.sections[0]
        section1.header.is_linked_to_previous = False
        section1.header.paragraphs[0].text = "Header text here"
        section1.footer.is_linked_to_previous = False
        section1.footer.paragraphs[0].text = "Footer text here"

        # Second section left linked to the previous one (the default) - its
        # header/footer must NOT be re-emitted.
        doc.add_section(WD_SECTION.NEW_PAGE)
        doc.add_paragraph("Section 2 body")

        buf_path = tmp_path / "sections.docx"
        doc.save(str(buf_path))

        # Inject a minimal footnotes.xml part - python-docx has no API to
        # add footnotes, so this is done directly at the OOXML level,
        # mirroring how a real Word-authored footnote is packaged.
        with zipfile.ZipFile(str(buf_path)) as zf:
            contents = {n: zf.read(n) for n in zf.namelist()}

        content_types = contents["[Content_Types].xml"].decode("utf-8")
        doc_rels = contents["word/_rels/document.xml.rels"].decode("utf-8")

        footnotes_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            '<w:footnote w:type="separator" w:id="-1">'
            "<w:p><w:r><w:separator/></w:r></w:p></w:footnote>"
            '<w:footnote w:type="continuationSeparator" w:id="0">'
            "<w:p><w:r><w:continuationSeparator/></w:r></w:p></w:footnote>"
            '<w:footnote w:id="1"><w:p><w:r><w:t>Footnote body text</w:t></w:r></w:p></w:footnote>'
            "</w:footnotes>"
        )
        content_types = content_types.replace(
            "</Types>",
            '<Override PartName="/word/footnotes.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml'
            '.footnotes+xml"/></Types>',
        )
        doc_rels = doc_rels.replace(
            "</Relationships>",
            '<Relationship Id="rIdFootnotes" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/footnotes" '
            'Target="footnotes.xml"/></Relationships>',
        )
        contents["[Content_Types].xml"] = content_types.encode("utf-8")
        contents["word/_rels/document.xml.rels"] = doc_rels.encode("utf-8")
        contents["word/footnotes.xml"] = footnotes_xml.encode("utf-8")

        with zipfile.ZipFile(str(buf_path), "w", zipfile.ZIP_DEFLATED) as zf:
            for name, data in contents.items():
                zf.writestr(name, data)

        result = list(self.ext.extract_stream(buf_path, MAX_SIZE))

        assert result[0] == "Body paragraph text"
        assert result[1] == "Section 2 body"
        assert "[Header] Header text here" in result
        assert "[Footer] Footer text here" in result
        assert "[Footnote] Footnote body text" in result
        # Linked second-section header must not duplicate section 1's header.
        assert result.count("[Header] Header text here") == 1
        # No boilerplate separator/continuationSeparator entries leaked.
        assert not any("separator" in r.lower() for r in result if "Footnote" in r)

    def test_docx_zip_bomb_prevention(self, tmp_path):
        fake_docx = tmp_path / "zip_bomb.docx"
        fake_docx.touch()

        mock_info = MagicMock()
        mock_info.file_size = 101 * 1024 * 1024
        mock_info.compress_size = 1024

        mock_zf = MagicMock()
        mock_zf.__enter__.return_value = mock_zf
        mock_zf.infolist.return_value = [mock_info]

        mock_docx = MagicMock()

        with (
            patch("zipfile.ZipFile", return_value=mock_zf),
            patch.dict("sys.modules", {"docx": mock_docx}),
        ):
            result = self.ext.extract(fake_docx, MAX_SIZE)

        assert result == ""
        mock_docx.Document.assert_not_called()


# ── XLSX (mocked openpyxl) ────────────────────────────────────────────────────


class TestXlsxExtractor:
    def setup_method(self):
        self.ext = XlsxExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "sheet.xlsx")
        assert self.ext.can_handle(tmp_path / "sheet.xls")
        assert not self.ext.can_handle(tmp_path / "sheet.csv")

    def test_extract_success(self, tmp_path):
        fake_path = tmp_path / "test.xlsx"
        fake_path.touch()
        mock_openpyxl = MagicMock()
        mock_ws = MagicMock()
        mock_ws.title = "Sheet1"
        mock_ws.iter_rows.return_value = iter(
            [
                ("Header", "Value"),
                ("Row1", 42),
            ]
        )
        mock_wb = MagicMock()
        mock_wb.worksheets = [mock_ws]
        mock_openpyxl.load_workbook.return_value = mock_wb
        with patch.dict("sys.modules", {"openpyxl": mock_openpyxl}):
            result = self.ext.extract(fake_path, MAX_SIZE)
        assert "Sheet1" in result
        assert "Header" in result
        assert "Row1" in result

    def test_extract_xlsx_size_truncation(self, tmp_path):
        fake_path = tmp_path / "test.xlsx"
        fake_path.touch()
        mock_openpyxl = MagicMock()
        mock_ws = MagicMock()
        mock_ws.title = "Sheet1"
        mock_ws.iter_rows.return_value = iter(
            [
                ("HeaderLongValue1234567890", "Value"),
            ]
        )
        mock_wb = MagicMock()
        mock_wb.worksheets = [mock_ws]
        mock_openpyxl.load_workbook.return_value = mock_wb
        with patch.dict("sys.modules", {"openpyxl": mock_openpyxl}):
            result = self.ext.extract(fake_path, 10)
        # Should truncate or exit loop early when size exceeded
        assert len(result) <= 30

    def test_extract_legacy_xls(self, tmp_path):
        fake_path = tmp_path / "test.xls"
        fake_path.touch()

        mock_xlrd = MagicMock()
        mock_wb = MagicMock()
        mock_sheet = MagicMock()
        mock_sheet.name = "LegacySheet"
        mock_sheet.nrows = 2

        c1 = MagicMock()
        c1.value = "H1"
        c2 = MagicMock()
        c2.value = "V1"
        c3 = MagicMock()
        c3.value = "H2"
        c4 = MagicMock()
        c4.value = "V2"

        mock_sheet.row.side_effect = [[c1, c2], [c3, c4]]
        mock_wb.sheets.return_value = [mock_sheet]
        mock_xlrd.open_workbook.return_value = mock_wb

        with patch.dict("sys.modules", {"xlrd": mock_xlrd}):
            result = self.ext.extract(fake_path, MAX_SIZE)
        assert "LegacySheet" in result
        assert "H1 | V1" in result
        assert "H2 | V2" in result

    def test_extract_legacy_xls_size_truncation(self, tmp_path):
        fake_path = tmp_path / "test.xls"
        fake_path.touch()

        mock_xlrd = MagicMock()
        mock_wb = MagicMock()
        mock_sheet = MagicMock()
        mock_sheet.name = "S"
        mock_sheet.nrows = 10
        c = MagicMock()
        c.value = "SomeLongTextValue"
        mock_sheet.row.return_value = [c]
        mock_wb.sheets.return_value = [mock_sheet]
        mock_xlrd.open_workbook.return_value = mock_wb

        with patch.dict("sys.modules", {"xlrd": mock_xlrd}):
            result = self.ext.extract(fake_path, 5)
        # Should stop processing sheets
        assert len(result) < 50

    def test_missing_file_returns_empty(self, tmp_path):
        result = self.ext.extract(tmp_path / "missing.xlsx", MAX_SIZE)
        assert result == ""


# -- PPTX (real decks; the extractor reads OOXML directly) ---------------------
#
# These were mocked against python-pptx, which the extractor no longer uses at
# all - it reads the zip's XML parts so peak memory is bounded by one slide
# rather than by an lxml tree for the whole package. Mocks could not have caught
# a regression in that rewrite, so every case below builds a real .pptx.
# Equivalence with the old python-pptx output is asserted separately, in
# tests/test_pptx_extractor_stream.py.


def _deck(path, *, title=None, body=None, table=None, notes=None):
    """Build a real single-slide .pptx and return its path."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if title is not None:
        slide.shapes.title.text = title
    if body is not None:
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = body
    if table is not None:
        rows, cols = len(table), len(table[0])
        shape = slide.shapes.add_table(
            rows, cols, Inches(1), Inches(3), Inches(6), Inches(0.8 * rows)
        )
        for r, row in enumerate(table):
            for c, val in enumerate(row):
                shape.table.cell(r, c).text = val
    if notes is not None:
        slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))
    return path


class TestPptxExtractor:
    def setup_method(self):
        self.ext = PptxExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "slides.pptx")
        assert not self.ext.can_handle(tmp_path / "slides.xlsx")

    def test_extract_success(self, tmp_path):
        path = _deck(tmp_path / "ok.pptx", title="Hello", body="Slide text content")
        result = self.ext.extract(path, MAX_SIZE)
        assert "Hello" in result
        assert "Slide text content" in result
        assert "--- Slide 1 ---" in result

    def test_missing_file_returns_empty(self, tmp_path):
        assert self.ext.extract(tmp_path / "missing.pptx", MAX_SIZE) == ""

    def test_pptx_zip_bomb_prevention(self, tmp_path):
        """The guard reads the central directory and must refuse before any
        part is decompressed."""
        fake_pptx = tmp_path / "zip_bomb.pptx"
        fake_pptx.touch()

        mock_info = MagicMock()
        mock_info.file_size = 101 * 1024 * 1024
        mock_info.compress_size = 1024
        mock_zf = MagicMock()
        mock_zf.__enter__.return_value = mock_zf
        mock_zf.infolist.return_value = [mock_info]

        with patch("zipfile.ZipFile", return_value=mock_zf):
            result = self.ext.extract(fake_pptx, MAX_SIZE)

        assert result == ""
        mock_zf.open.assert_not_called()

    def test_extract_includes_table_cells_after_title(self, tmp_path):
        """P0-3: GraphicFrame tables were silently dropped. Cell text must
        appear, and must come AFTER the regular shape text - the summarizer
        takes the first line after each slide marker as that slide's title,
        so table content can't displace it."""
        path = _deck(
            tmp_path / "table.pptx",
            title="Q4 Overview",
            table=[["Region", "Revenue"], ["EMEA", "12"]],
        )
        out = list(self.ext.extract_stream(path, MAX_SIZE))

        assert out[0] == "--- Slide 1 ---"
        assert out[1] == "Q4 Overview", "table text displaced the slide title"
        joined = "\n".join(out)
        assert "Region | Revenue" in joined
        assert "EMEA | 12" in joined

    def test_extract_includes_speaker_notes(self, tmp_path):
        path = _deck(tmp_path / "notes.pptx", title="Deck", notes="Remember the caveat")
        joined = "\n".join(self.ext.extract_stream(path, MAX_SIZE))
        assert "[Notes] Remember the caveat" in joined

    def test_notes_exclude_the_slide_number_placeholder(self, tmp_path):
        """A notes slide also carries a slide-number placeholder and a
        thumbnail of the slide. Taking every text frame on the part appended a
        stray page number to every note - caught by the differential test
        against real decks, not by review."""
        path = _deck(tmp_path / "notesonly.pptx", title="Deck", notes="Just this")
        notes = [ln for ln in self.ext.extract_stream(path, MAX_SIZE) if ln.startswith("[Notes]")]
        assert notes == ["[Notes] Just this"]

    def test_slide_without_notes_emits_none(self, tmp_path):
        path = _deck(tmp_path / "plain.pptx", title="Deck", body="text")
        assert not [
            ln for ln in self.ext.extract_stream(path, MAX_SIZE) if ln.startswith("[Notes]")
        ]

    def test_corrupt_file_does_not_raise(self, tmp_path):
        bad = tmp_path / "corrupt.pptx"
        bad.write_bytes(b"not a zip at all")
        assert self.ext.extract(bad, MAX_SIZE) == ""

    def test_slide_title_survives_as_first_line_for_summarizer(self, tmp_path):
        from app.indexing.summarizer import _summarize_doc_text

        path = _deck(
            tmp_path / "sum.pptx",
            title="Q4 Overview",
            table=[["Region", "Revenue"]],
            notes="internal only",
        )
        result = self.ext.extract(path, MAX_SIZE)
        summary = _summarize_doc_text(result, 300)
        assert "Q4 Overview" in summary


# ── EPUB ──────────────────────────────────────────────────────────────────────


class TestEpubExtractor:
    def setup_method(self):
        self.ext = EpubExtractor()

    def test_can_handle(self, tmp_path):
        assert self.ext.can_handle(tmp_path / "book.epub")
        assert not self.ext.can_handle(tmp_path / "book.pdf")

    def test_missing_file_returns_empty_or_error(self, tmp_path):
        result = self.ext.extract(tmp_path / "missing.epub", MAX_SIZE)
        assert result == ""

    def test_epub_extraction_success(self, tmp_path):
        fake_epub = tmp_path / "test.epub"

        # We can write a real zip file to test real EPUB structure
        import zipfile

        with zipfile.ZipFile(fake_epub, "w") as zf:
            zf.writestr("mimetype", "application/epub+zip")
            # Subfolder OEBPS/ or any folder
            zf.writestr(
                "OEBPS/ch1.xhtml",
                "<html><body><h1>Introduction</h1><p>Welcome to &amp; standard text. This is a longer text paragraph that contains enough characters to pass the fifty character check in the extractor.</p></body></html>",
            )
            zf.writestr(
                "OEBPS/ch2.html",
                "<html><body><p>This is second page. It has a lot of additional characters in order to exceed fifty characters as well.</p></body></html>",
            )
            # Ignored folder
            zf.writestr("__MACOSX/ch1.xhtml", "ignored")
            # Non-html file suffix
            zf.writestr("OEBPS/image.png", "binarydata")

        result = self.ext.extract(fake_epub, MAX_SIZE)
        assert "Introduction" in result
        assert "Welcome to standard text." in result
        assert "This is second page." in result
        assert "ignored" not in result

    def test_epub_strips_inline_script_and_style_content(self, tmp_path):
        """_HTML_TAG_RE alone strips tags but leaves enclosed text, so a
        <script> body or <style> ruleset would otherwise land in the index
        as prose."""
        fake_epub = tmp_path / "script_style.epub"

        with zipfile.ZipFile(fake_epub, "w") as zf:
            zf.writestr("mimetype", "application/epub+zip")
            zf.writestr(
                "OEBPS/ch1.xhtml",
                "<html><head><style>body { color: red; font-size: 12px; } "
                ".warning-class-name { display: none; }</style>"
                "<script>function trackEvent() { alert('leak-marker-token'); "
                "console.log('another leak'); }</script></head>"
                "<body><p>Legitimate chapter prose that is long enough to "
                "clear the fifty character minimum threshold check easily.</p>"
                "</body></html>",
            )

        result = self.ext.extract(fake_epub, MAX_SIZE)
        assert "Legitimate chapter prose" in result
        assert "leak-marker-token" not in result
        assert "trackEvent" not in result
        assert "color: red" not in result
        assert "warning-class-name" not in result

    def test_epub_zip_bomb_prevention(self, tmp_path):
        fake_epub = tmp_path / "zip_bomb.epub"
        fake_epub.touch()

        mock_zf = MagicMock()
        mock_zf.__enter__.return_value = mock_zf
        mock_zf.namelist.return_value = ["ch1.xhtml"]
        mock_zf.getinfo.return_value.file_size = 101 * 1024 * 1024

        with patch("zipfile.ZipFile", return_value=mock_zf):
            result = self.ext.extract(fake_epub, MAX_SIZE)
        # Oversized entries are rejected before they are decompressed.
        assert result == ""
        mock_zf.open.assert_not_called()

    def test_epub_uses_opf_spine_order_not_alphabetical(self, tmp_path):
        """P0-3: alphabetically 'content_a.xhtml' < 'content_z.xhtml', but the
        spine lists content_z first - a real reader must follow the spine,
        not the filename. This EPUB is built so alphabetical order would
        read the chapters backwards."""
        fake_epub = tmp_path / "spine_order.epub"

        import zipfile

        container_xml = """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>"""
        content_opf = """<?xml version="1.0"?>
<package xmlns="http://www.idpf.org/2007/opf" version="2.0">
  <manifest>
    <item id="chap-z" href="content_z.xhtml" media-type="application/xhtml+xml"/>
    <item id="chap-a" href="content_a.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="chap-z"/>
    <itemref idref="chap-a"/>
  </spine>
</package>"""

        with zipfile.ZipFile(fake_epub, "w") as zf:
            zf.writestr("mimetype", "application/epub+zip")
            zf.writestr("META-INF/container.xml", container_xml)
            zf.writestr("OEBPS/content.opf", content_opf)
            zf.writestr(
                "OEBPS/content_z.xhtml",
                "<html><body><p>This is the FIRST chapter per the spine order, "
                "padded with enough characters to clear the fifty character floor.</p>"
                "</body></html>",
            )
            zf.writestr(
                "OEBPS/content_a.xhtml",
                "<html><body><p>This is the SECOND chapter per the spine order, "
                "padded with enough characters to clear the fifty character floor.</p>"
                "</body></html>",
            )

        result = self.ext.extract(fake_epub, MAX_SIZE)
        assert "FIRST chapter" in result
        assert "SECOND chapter" in result
        assert result.index("FIRST chapter") < result.index("SECOND chapter")

    def test_epub_falls_back_to_alphabetical_without_container_xml(self, tmp_path, caplog):
        """No META-INF/container.xml at all - must not crash, must fall back,
        and must log the fallback at INFO so a missing spine is diagnosable."""
        fake_epub = tmp_path / "no_container.epub"

        import zipfile

        with zipfile.ZipFile(fake_epub, "w") as zf:
            zf.writestr(
                "chapter_b.xhtml",
                "<html><body><p>Content of chapter B, long enough to clear the "
                "fifty character minimum fragment length floor in the extractor.</p></body></html>",
            )

        with caplog.at_level("INFO", logger="app.indexing.extractors.epub_extractor"):
            result = self.ext.extract(fake_epub, MAX_SIZE)

        assert "Content of chapter B" in result
        assert any("falling back to alphabetical" in r.message for r in caplog.records)

    def test_epub_falls_back_when_spine_references_missing_entries(self, tmp_path):
        """A spine that points at hrefs not actually present in the zip
        (malformed/incomplete EPUB) must fall back rather than yield nothing."""
        fake_epub = tmp_path / "broken_spine.epub"

        import zipfile

        container_xml = """<?xml version="1.0"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles>
    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>
  </rootfiles>
</container>"""
        content_opf = """<?xml version="1.0"?>
<package xmlns="http://www.idpf.org/2007/opf" version="2.0">
  <manifest>
    <item id="ghost" href="does_not_exist.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine>
    <itemref idref="ghost"/>
  </spine>
</package>"""

        with zipfile.ZipFile(fake_epub, "w") as zf:
            zf.writestr("META-INF/container.xml", container_xml)
            zf.writestr("OEBPS/content.opf", content_opf)
            zf.writestr(
                "actual_content.xhtml",
                "<html><body><p>This is the only real chapter, long enough to clear "
                "the fifty character minimum fragment floor in the extractor.</p></body></html>",
            )

        result = self.ext.extract(fake_epub, MAX_SIZE)
        assert "only real chapter" in result


# ── EXTRACTORS registry ───────────────────────────────────────────────────────


class TestExtractorRegistry:
    def test_all_extractors_are_loaded(self):
        assert len(EXTRACTORS) == 7

    def test_each_extractor_has_can_handle_and_extract(self):
        for ext in EXTRACTORS:
            assert callable(getattr(ext, "can_handle", None))
            assert callable(getattr(ext, "extract", None))

    def test_only_one_extractor_handles_each_type(self, tmp_path):
        types = [".pdf", ".docx", ".xlsx", ".pptx", ".epub", ".csv", ".json"]
        for t in types:
            handlers = [e for e in EXTRACTORS if e.can_handle(tmp_path / f"file{t}")]
            assert len(handlers) == 1, f"Expected exactly 1 handler for {t}, got {len(handlers)}"
